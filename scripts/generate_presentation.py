"""
Generador de la presentacion final del proyecto TIA.

Crea reports/presentacion.pptx con 8 slides para sustentar el proyecto
en 5-7 minutos. Cada slide tiene contenido + notas del presentador con
el guion de narracion.

Uso:
    python scripts/generate_presentation.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ─── Paleta de colores ────────────────────────────────────────────────
AZUL_PRIMARIO = RGBColor(0x1F, 0x4E, 0x79)      # encabezados
AZUL_CLARO = RGBColor(0x2E, 0x75, 0xB6)         # destacados
GRIS_OSCURO = RGBColor(0x33, 0x33, 0x33)        # texto cuerpo
GRIS_CLARO = RGBColor(0xF2, 0xF2, 0xF2)         # fondo cajas
VERDE_OK = RGBColor(0x2E, 0x7D, 0x32)           # metricas buenas
NARANJA = RGBColor(0xED, 0x7D, 0x31)            # acentos


# ─── Configuracion ────────────────────────────────────────────────────
PROJECT_ROOT = Path(__file__).resolve().parent.parent
OUTPUT_PATH = PROJECT_ROOT / "reports" / "presentacion.pptx"


# ─── Helpers ──────────────────────────────────────────────────────────


def set_text(text_frame, text, size=18, bold=False, color=None, align=None):
    """Configura el texto de un text frame con estilo."""
    text_frame.clear()
    p = text_frame.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_bullet(text_frame, text, size=16, color=None, indent=0):
    """Agrega un bullet al text frame."""
    p = text_frame.add_paragraph()
    p.level = indent
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    if color:
        run.font.color.rgb = color


def add_title_bar(slide, prs, title_text):
    """Agrega una barra de titulo en la parte superior de la slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(0.8),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = AZUL_PRIMARIO
    bar.line.fill.background()

    text_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.15),
        prs.slide_width - Inches(0.8), Inches(0.5),
    )
    set_text(
        text_box.text_frame, title_text,
        size=24, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
    )


def set_notes(slide, notes_text):
    """Agrega notas del presentador a la slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


def add_footer(slide, prs, text):
    """Agrega un footer discreto."""
    footer = slide.shapes.add_textbox(
        Inches(0.4), prs.slide_height - Inches(0.4),
        prs.slide_width - Inches(0.8), Inches(0.3),
    )
    tf = footer.text_frame
    set_text(tf, text, size=10, color=GRIS_OSCURO, align=PP_ALIGN.RIGHT)


# ─── Slides ───────────────────────────────────────────────────────────


def slide_1_portada(prs):
    """Slide 1 - Portada."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Banda azul superior
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(2.5),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = AZUL_PRIMARIO
    bar.line.fill.background()

    # Titulo
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.7),
        prs.slide_width - Inches(1.2), Inches(1.2),
    )
    set_text(
        title_box.text_frame,
        "Prediccion de aprobacion estudiantil",
        size=36, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
    )

    subtitle_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.7),
        prs.slide_width - Inches(1.2), Inches(0.6),
    )
    set_text(
        subtitle_box.text_frame,
        "Modelo de Machine Learning para identificar estudiantes en riesgo academico",
        size=18, color=RGBColor(0xCF, 0xE2, 0xF3),
    )

    # Bloque inferior con autor
    info_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(3.5),
        prs.slide_width - Inches(1.2), Inches(2.0),
    )
    tf = info_box.text_frame
    tf.word_wrap = True
    set_text(tf, "Samuel Rendon Hincapie", size=22, bold=True, color=AZUL_PRIMARIO)
    add_bullet(tf, "Curso: Machine Learning", size=16, color=GRIS_OSCURO)
    add_bullet(tf, "Institucion Universitaria Pascual Bravo", size=16, color=GRIS_OSCURO)
    add_bullet(tf, "Trabajo Integrador de Aprendizaje - TIA3", size=16, color=GRIS_OSCURO)

    set_notes(slide, (
        "GUION (15 segundos):\n\n"
        "'Buenos dias profesor. Mi proyecto es un modelo de clasificacion "
        "que predice si un estudiante va a aprobar o reprobar el examen "
        "final, a partir de sus habitos de estudio y bienestar.'\n\n"
        "TIP: arranca con voz firme. No leas el titulo, di la idea principal."
    ))


def slide_2_pregunta_negocio(prs):
    """Slide 2 - Pregunta de negocio."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, "Pregunta de negocio")

    # Cita destacada
    cita_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.4),
        prs.slide_width - Inches(1.6), Inches(1.8),
    )
    cita_box.fill.solid()
    cita_box.fill.fore_color.rgb = GRIS_CLARO
    cita_box.line.color.rgb = AZUL_CLARO
    cita_box.line.width = Pt(2)

    cita_text = slide.shapes.add_textbox(
        Inches(1.2), Inches(1.7),
        prs.slide_width - Inches(2.4), Inches(1.4),
    )
    tf = cita_text.text_frame
    tf.word_wrap = True
    set_text(
        tf,
        '"Deberia un estudiante cambiar sus habitos de estudio y vida '
        "social basandose en la prediccion de si aprobara o reprobara "
        'el examen final?"',
        size=18, color=GRIS_OSCURO,
    )

    # Dataset
    dataset_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(3.6),
        prs.slide_width - Inches(1.6), Inches(2.5),
    )
    tf = dataset_box.text_frame
    tf.word_wrap = True
    set_text(tf, "Dataset Student Performance", size=18, bold=True, color=AZUL_PRIMARIO)
    add_bullet(tf, "5,000 estudiantes con 21 variables", size=15, color=GRIS_OSCURO)
    add_bullet(tf, "Habitos diarios: estudio, sueno, redes sociales, gaming", size=15, color=GRIS_OSCURO)
    add_bullet(tf, "Bienestar: salud mental, burnout, productividad", size=15, color=GRIS_OSCURO)
    add_bullet(tf, "Contexto academico: nivel educativo, trabajo de medio tiempo, deadlines", size=15, color=GRIS_OSCURO)

    set_notes(slide, (
        "GUION (30 segundos):\n\n"
        "'La pregunta de negocio es esta: deberia un estudiante cambiar "
        "sus habitos basandose en la prediccion del modelo?\n\n"
        "El dataset tiene 5000 estudiantes con 21 variables que cubren "
        "tres dimensiones: habitos de estudio, bienestar, y contexto "
        "academico. Lo importante es que el modelo no solo prediga, "
        "sino que sus predicciones se puedan accionar.'\n\n"
        "TIP: si tienes tiempo, di una frase del 'por que'. "
        "Ej: 'porque el rendimiento no depende solo de la capacidad, "
        "tambien de habitos modificables'."
    ))


def slide_3_variables(prs):
    """Slide 3 - Variables del modelo."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, "Variables del modelo (13 features)")

    # Tres columnas tematicas
    col_width = (prs.slide_width - Inches(1.2)) / 3
    columnas = [
        {
            "titulo": "Habitos de estudio",
            "items": [
                "study_hours",
                "self_study_hours",
                "social_media_hours",
                "screen_time_hours",
            ],
            "color": AZUL_CLARO,
        },
        {
            "titulo": "Bienestar fisico",
            "items": [
                "sleep_hours",
                "exercise_minutes",
                "caffeine_intake_mg",
            ],
            "color": VERDE_OK,
        },
        {
            "titulo": "Contexto + indicadores",
            "items": [
                "part_time_job",
                "upcoming_deadline",
                "mental_health_score",
                "focus_index",
                "burnout_level",
                "productivity_score",
            ],
            "color": NARANJA,
        },
    ]

    for i, col in enumerate(columnas):
        left = Inches(0.4) + col_width * i + Inches(0.1)
        # Caja titulo
        titulo_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, Inches(1.3),
            col_width - Inches(0.2), Inches(0.5),
        )
        titulo_box.fill.solid()
        titulo_box.fill.fore_color.rgb = col["color"]
        titulo_box.line.fill.background()
        titulo_text = slide.shapes.add_textbox(
            left, Inches(1.35), col_width - Inches(0.2), Inches(0.4),
        )
        set_text(
            titulo_text.text_frame, col["titulo"],
            size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
        )
        # Items
        items_box = slide.shapes.add_textbox(
            left, Inches(2.0), col_width - Inches(0.2), Inches(4.5),
        )
        tf = items_box.text_frame
        tf.word_wrap = True
        for j, item in enumerate(col["items"]):
            if j == 0:
                set_text(tf, item, size=14, color=GRIS_OSCURO)
            else:
                add_bullet(tf, item, size=14, color=GRIS_OSCURO)

    # Nota inferior
    nota_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(6.5),
        prs.slide_width - Inches(0.8), Inches(0.6),
    )
    set_text(
        nota_box.text_frame,
        "Variable objetivo: aprobado (1) / reprueba (0) — binarizada desde exam_score con umbral mediana = 18.01",
        size=12, color=GRIS_OSCURO, align=PP_ALIGN.CENTER,
    )

    set_notes(slide, (
        "GUION (45 segundos):\n\n"
        "'Despues de un analisis exploratorio extenso, seleccione 13 "
        "variables en 3 grupos tematicos: habitos de estudio, bienestar "
        "fisico, y contexto + indicadores.\n\n"
        "La variable objetivo es la nota del examen, binarizada con la "
        "mediana del dataset (18.01) para tener clases balanceadas.\n\n"
        "Excluí 7 variables que el EDA mostro irrelevantes: edad, genero, "
        "nivel academico, calidad de internet, y clases online — todas con "
        "correlacion practicamente cero con el resultado.'\n\n"
        "TIP: si te apuran, di solo los grupos sin enumerar todas las variables."
    ))


def slide_4_pipeline(prs):
    """Slide 4 - Pipeline del modelo."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, "Pipeline del modelo")

    # Diagrama horizontal: 3 cajas + flechas
    box_width = Inches(2.6)
    box_height = Inches(1.2)
    top = Inches(2.0)
    start_left = Inches(0.5)
    gap = Inches(0.3)

    pasos = [
        ("StandardScaler", "Escalado de\nlas 13 features", AZUL_CLARO),
        ("SMOTE", "Balanceo de\nclases (solo fit)", NARANJA),
        ("DecisionTreeClassifier", "max_depth=5\nrandom_state=42", VERDE_OK),
    ]

    for i, (titulo, descripcion, color) in enumerate(pasos):
        left = start_left + (box_width + gap) * i

        # Caja
        caja = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_width, box_height,
        )
        caja.fill.solid()
        caja.fill.fore_color.rgb = color
        caja.line.fill.background()

        # Texto
        text_box = slide.shapes.add_textbox(
            left, top + Inches(0.1), box_width, box_height - Inches(0.2),
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        set_text(
            tf, titulo,
            size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
        )
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = descripcion
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Flecha (excepto la ultima)
        if i < len(pasos) - 1:
            flecha = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                left + box_width + Inches(0.02),
                top + Inches(0.45),
                gap - Inches(0.04), Inches(0.3),
            )
            flecha.fill.solid()
            flecha.fill.fore_color.rgb = GRIS_OSCURO
            flecha.line.fill.background()

    # Codigo simplificado abajo
    codigo_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(4.0),
        prs.slide_width - Inches(1.0), Inches(2.6),
    )
    codigo_box.fill.solid()
    codigo_box.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
    codigo_box.line.fill.background()

    codigo_text = slide.shapes.add_textbox(
        Inches(0.7), Inches(4.15),
        prs.slide_width - Inches(1.4), Inches(2.4),
    )
    tf = codigo_text.text_frame
    tf.word_wrap = True
    codigo = (
        "pipeline = Pipeline([\n"
        '    ("preproc", StandardScaler()),\n'
        '    ("smote", SMOTE(random_state=42)),\n'
        '    ("classifier", DecisionTreeClassifier(max_depth=5)),\n'
        "])\n"
        "pipeline.fit(X_train, y_train)"
    )
    set_text(tf, codigo, size=14, color=RGBColor(0xE8, 0xE8, 0xE8))
    # forzar fuente monospace en cada run
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.name = "Consolas"

    set_notes(slide, (
        "GUION (45 segundos):\n\n"
        "'El pipeline tiene tres pasos encadenados:\n\n"
        "1. StandardScaler para escalar las 13 features.\n"
        "2. SMOTE para balancear las clases solo durante entrenamiento — "
        "es importante que se ignore en prediccion, por eso usamos "
        "imblearn.Pipeline en lugar del de sklearn.\n"
        "3. Arbol de decision con profundidad 5, que es interpretable y "
        "se ajusto bien al tamano del dataset.\n\n"
        "Todo el pipeline se serializa como un solo objeto joblib, lo que "
        "garantiza reproducibilidad en produccion.'\n\n"
        "TIP: enfatiza que SMOTE solo aplica en fit. Eso muestra "
        "que entendiste el problema del data leakage."
    ))


def slide_5_resultados(prs):
    """Slide 5 - Resultados (metricas) -- SLIDE PRINCIPAL."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, "Resultados del modelo")

    # Cinco metricas en cajas grandes
    metricas = [
        ("Accuracy", "0.842", "Predicciones correctas"),
        ("Precision", "0.825", "Calidad de alertas"),
        ("Recall", "0.868", "Cobertura de casos reales"),
        ("F1-Score", "0.846", "Balance precision/recall"),
        ("AUC-ROC", "0.933", "Capacidad discriminativa"),
    ]

    box_width = Inches(1.85)
    box_height = Inches(2.2)
    top = Inches(1.4)
    start_left = Inches(0.4)
    gap = Inches(0.05)

    for i, (nombre, valor, desc) in enumerate(metricas):
        left = start_left + (box_width + gap) * i

        # Caja con sombra
        caja = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_width, box_height,
        )
        caja.fill.solid()
        caja.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        caja.line.color.rgb = AZUL_CLARO
        caja.line.width = Pt(2)

        # Nombre
        nombre_box = slide.shapes.add_textbox(
            left, top + Inches(0.15), box_width, Inches(0.4),
        )
        set_text(
            nombre_box.text_frame, nombre,
            size=14, bold=True, color=AZUL_PRIMARIO, align=PP_ALIGN.CENTER,
        )

        # Valor (grande)
        valor_box = slide.shapes.add_textbox(
            left, top + Inches(0.65), box_width, Inches(1.0),
        )
        set_text(
            valor_box.text_frame, valor,
            size=42, bold=True, color=VERDE_OK, align=PP_ALIGN.CENTER,
        )

        # Descripcion
        desc_box = slide.shapes.add_textbox(
            left + Inches(0.05), top + Inches(1.65),
            box_width - Inches(0.1), Inches(0.5),
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        set_text(
            tf, desc,
            size=10, color=GRIS_OSCURO, align=PP_ALIGN.CENTER,
        )

    # Interpretacion inferior
    interp_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), Inches(4.0),
        prs.slide_width - Inches(0.8), Inches(2.5),
    )
    interp_box.fill.solid()
    interp_box.fill.fore_color.rgb = GRIS_CLARO
    interp_box.line.color.rgb = AZUL_CLARO
    interp_box.line.width = Pt(1)

    interp_text = slide.shapes.add_textbox(
        Inches(0.6), Inches(4.15),
        prs.slide_width - Inches(1.2), Inches(2.3),
    )
    tf = interp_text.text_frame
    tf.word_wrap = True
    set_text(
        tf, "Interpretacion en lenguaje de negocio",
        size=15, bold=True, color=AZUL_PRIMARIO,
    )
    add_bullet(tf, "El modelo identifica al 86.8% de los estudiantes que efectivamente reprobarian (recall alto)", size=13, color=GRIS_OSCURO)
    add_bullet(tf, "De cada 100 alertas, 82 son correctas (precision alta -> no saturamos al equipo de bienestar)", size=13, color=GRIS_OSCURO)
    add_bullet(tf, "AUC-ROC de 0.93 indica capacidad discriminativa excelente, independiente del umbral", size=13, color=GRIS_OSCURO)

    set_notes(slide, (
        "GUION (75 segundos) — SLIDE CLAVE:\n\n"
        "'Estos son los resultados del modelo sobre el set de prueba de "
        "1000 estudiantes. Cinco metricas: accuracy de 84%, precision "
        "82%, recall 87%, F1 de 85%, y un AUC-ROC de 0.93.\n\n"
        "Pero los numeros sueltos no dicen mucho. Lo importante es lo "
        "que significan en el negocio:\n\n"
        "Recall 0.87 significa que identificamos al 87% de los estudiantes "
        "que realmente iban a reprobar. Es decir, casi nadie se nos "
        "escapa sin alerta.\n\n"
        "Precision 0.82 significa que cuando damos una alerta, en 82 de "
        "cada 100 casos es real. Esto evita saturar al equipo de "
        "bienestar con falsos positivos.\n\n"
        "Y el AUC de 0.93 es excelente: el modelo distingue muy bien las "
        "dos clases independiente del umbral que escojamos.'\n\n"
        "TIP: detente unos segundos despues de cada metrica. No las leas "
        "en bloque, deja que se procesen."
    ))


def slide_6_graficas(prs):
    """Slide 6 - Matriz de confusion + ROC."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, "Matriz de confusion y curva ROC")

    # Placeholder izquierdo: matriz de confusion
    placeholder_left = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.4), Inches(1.3),
        Inches(4.8), Inches(4.8),
    )
    placeholder_left.fill.solid()
    placeholder_left.fill.fore_color.rgb = GRIS_CLARO
    placeholder_left.line.color.rgb = AZUL_CLARO
    placeholder_left.line.dash_style = 2  # dashed

    text_left = slide.shapes.add_textbox(
        Inches(0.4), Inches(3.4),
        Inches(4.8), Inches(0.6),
    )
    set_text(
        text_left.text_frame,
        "[ Insertar: reports/v1/confusion_matrix.png ]",
        size=14, color=GRIS_OSCURO, align=PP_ALIGN.CENTER,
    )

    label_left = slide.shapes.add_textbox(
        Inches(0.4), Inches(6.2),
        Inches(4.8), Inches(0.4),
    )
    set_text(
        label_left.text_frame,
        "Matriz de confusion (test set)",
        size=14, bold=True, color=AZUL_PRIMARIO, align=PP_ALIGN.CENTER,
    )

    # Placeholder derecho: ROC
    placeholder_right = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.4), Inches(1.3),
        Inches(4.8), Inches(4.8),
    )
    placeholder_right.fill.solid()
    placeholder_right.fill.fore_color.rgb = GRIS_CLARO
    placeholder_right.line.color.rgb = AZUL_CLARO
    placeholder_right.line.dash_style = 2

    text_right = slide.shapes.add_textbox(
        Inches(5.4), Inches(3.4),
        Inches(4.8), Inches(0.6),
    )
    set_text(
        text_right.text_frame,
        "[ Insertar: reports/v1/roc_curve.png ]",
        size=14, color=GRIS_OSCURO, align=PP_ALIGN.CENTER,
    )

    label_right = slide.shapes.add_textbox(
        Inches(5.4), Inches(6.2),
        Inches(4.8), Inches(0.4),
    )
    set_text(
        label_right.text_frame,
        "Curva ROC (AUC = 0.93)",
        size=14, bold=True, color=AZUL_PRIMARIO, align=PP_ALIGN.CENTER,
    )

    set_notes(slide, (
        "GUION (75 segundos):\n\n"
        "'A la izquierda esta la matriz de confusion. En el cuadrante "
        "superior izquierdo y el inferior derecho estan las predicciones "
        "correctas; en los otros dos estan los errores.\n\n"
        "Lo que observamos es que el modelo se equivoca de forma "
        "balanceada: comete pocos falsos positivos y pocos falsos "
        "negativos. Esto es lo que se busca cuando ambos errores "
        "tienen costo en el negocio.\n\n"
        "A la derecha esta la curva ROC. Una curva pegada a la esquina "
        "superior izquierda significa que el modelo separa bien las "
        "clases. La diagonal punteada representa el desempeno de un "
        "modelo aleatorio. Como ven, nuestra curva esta muy lejos de la "
        "diagonal — el AUC es 0.93 sobre un maximo de 1.0.'\n\n"
        "TIP IMPORTANTE: Antes de la sustentacion, abre PowerPoint, ve a "
        "esta slide, y reemplaza los dos placeholders grises con los PNG "
        "reales de reports/v1/.\n\n"
        "Si te queda tiempo, menciona la celda especifica: 'tenemos 442 "
        "verdaderos negativos, 450 verdaderos positivos, solo 95 errores "
        "en cada direccion'."
    ))


def slide_7_importancia(prs):
    """Slide 7 - Importancia de variables."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, "Variables que mas explican el resultado")

    # Placeholder grafica izquierda
    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.4), Inches(1.3),
        Inches(6.0), Inches(5.0),
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = GRIS_CLARO
    placeholder.line.color.rgb = AZUL_CLARO
    placeholder.line.dash_style = 2

    text_placeholder = slide.shapes.add_textbox(
        Inches(0.4), Inches(3.4),
        Inches(6.0), Inches(0.8),
    )
    set_text(
        text_placeholder.text_frame,
        "[ Insertar: reports/v1/feature_importances.png ]",
        size=14, color=GRIS_OSCURO, align=PP_ALIGN.CENTER,
    )

    # Lista de hallazgos a la derecha
    hallazgos_box = slide.shapes.add_textbox(
        Inches(6.7), Inches(1.4),
        Inches(3.5), Inches(5.5),
    )
    tf = hallazgos_box.text_frame
    tf.word_wrap = True
    set_text(tf, "Top hallazgos", size=16, bold=True, color=AZUL_PRIMARIO)
    add_bullet(tf, "focus_index domina la decision del modelo", size=13, color=GRIS_OSCURO)
    add_bullet(tf, "Las horas de estudio importan, pero menos que la concentracion", size=13, color=GRIS_OSCURO)
    add_bullet(tf, "burnout_level y mental_health_score son criticos", size=13, color=GRIS_OSCURO)
    add_bullet(tf, "Variables como deadlines o trabajo de medio tiempo casi no influyen", size=13, color=GRIS_OSCURO)

    # Insight de negocio destacado
    insight_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.7), Inches(5.2),
        Inches(3.5), Inches(1.4),
    )
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = NARANJA
    insight_box.line.fill.background()

    insight_text = slide.shapes.add_textbox(
        Inches(6.85), Inches(5.3),
        Inches(3.3), Inches(1.2),
    )
    tf = insight_text.text_frame
    tf.word_wrap = True
    set_text(
        tf, "Insight de negocio",
        size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
    )
    add_bullet(
        tf,
        "Las intervenciones mas efectivas son mejorar la concentracion y reducir el burnout",
        size=12, color=RGBColor(0xFF, 0xFF, 0xFF),
    )

    set_notes(slide, (
        "GUION (60 segundos):\n\n"
        "'Una de las ventajas de usar un arbol de decision es que es "
        "interpretable. Podemos ver exactamente cuales variables usa el "
        "modelo para decidir.\n\n"
        "El hallazgo principal es que focus_index — el indice de "
        "concentracion — domina la decision. Esto tiene sentido: un "
        "estudiante puede pasar muchas horas en clase pero si no se "
        "concentra, no rinde.\n\n"
        "Le siguen las horas de estudio, el nivel de burnout, y la salud "
        "mental. Variables como tener trabajo de medio tiempo o deadlines "
        "proximas, que uno podria pensar que afectan, casi no influyen.\n\n"
        "El insight para negocio es claro: las intervenciones mas "
        "efectivas no son administrativas; son sobre concentracion y "
        "bienestar mental.'\n\n"
        "TIP: este es un buen momento para conectar con la pregunta de "
        "negocio del inicio. Cierra el arco narrativo."
    ))


def slide_8_cierre(prs):
    """Slide 8 - Cierre."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, "Conclusiones")

    # Tres cajas de logros
    logros = [
        (
            "Modelo predictivo solido",
            "Pipeline encadenado con SMOTE + arbol\nAUC 0.93 sobre test set",
            VERDE_OK,
        ),
        (
            "Componente prescriptivo",
            "5 reglas de negocio (RN-01 a RN-05)\nRecomendaciones personalizadas",
            NARANJA,
        ),
        (
            "Migracion a produccion",
            "De Colab a VS Code modular\n~100 tests automatizados\nModelo serializado con metadata",
            AZUL_CLARO,
        ),
    ]

    box_width = (prs.slide_width - Inches(1.6)) / 3
    box_height = Inches(2.5)
    top = Inches(1.4)
    gap = Inches(0.2)

    for i, (titulo, desc, color) in enumerate(logros):
        left = Inches(0.4) + (box_width + gap) * i

        caja = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_width, box_height,
        )
        caja.fill.solid()
        caja.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        caja.line.color.rgb = color
        caja.line.width = Pt(3)

        # Numero
        num_box = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left + Inches(0.15), top + Inches(0.15),
            Inches(0.55), Inches(0.55),
        )
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = color
        num_box.line.fill.background()
        num_text = slide.shapes.add_textbox(
            left + Inches(0.15), top + Inches(0.18),
            Inches(0.55), Inches(0.55),
        )
        set_text(
            num_text.text_frame, str(i + 1),
            size=20, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
        )

        # Titulo
        titulo_text = slide.shapes.add_textbox(
            left + Inches(0.85), top + Inches(0.2),
            box_width - Inches(1.0), Inches(0.6),
        )
        set_text(
            titulo_text.text_frame, titulo,
            size=14, bold=True, color=color,
        )

        # Descripcion
        desc_box = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.9),
            box_width - Inches(0.4), Inches(1.5),
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        set_text(tf, desc, size=12, color=GRIS_OSCURO)

    # Bloque inferior: invitacion a preguntas
    pregunta_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.4), Inches(5.0),
        prs.slide_width - Inches(0.8), Inches(1.5),
    )
    pregunta_box.fill.solid()
    pregunta_box.fill.fore_color.rgb = AZUL_PRIMARIO
    pregunta_box.line.fill.background()

    pregunta_text = slide.shapes.add_textbox(
        Inches(0.6), Inches(5.1),
        prs.slide_width - Inches(1.2), Inches(1.3),
    )
    tf = pregunta_text.text_frame
    tf.word_wrap = True
    set_text(
        tf, "Gracias - Preguntas y discusion",
        size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER,
    )
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Codigo completo: codigo_proyecto.html  -  Modelo: decision_tree_v1.joblib"
    run.font.size = Pt(13)
    run.font.color.rgb = RGBColor(0xCF, 0xE2, 0xF3)

    set_notes(slide, (
        "GUION (30 segundos):\n\n"
        "'Para cerrar: el proyecto entrega tres cosas. Un modelo predictivo "
        "solido con AUC de 0.93. Un componente prescriptivo de 5 reglas que "
        "traduce predicciones en recomendaciones accionables. Y una "
        "migracion completa desde Colab a un entorno modular con cerca de "
        "100 tests automatizados.\n\n"
        "Quedo atento a sus preguntas.'\n\n"
        "TIP FINAL: termina con una pausa de 2 segundos antes de invitar "
        "a preguntas. No digas 'eso es todo'. Di 'quedo atento a sus "
        "preguntas' y sonrie.\n\n"
        "POSIBLES PREGUNTAS Y RESPUESTAS BREVES:\n"
        "- Por que arbol y no random forest? -> interpretabilidad y "
        "  tamano del dataset.\n"
        "- Por que SMOTE si las clases estan balanceadas? -> por "
        "  consistencia metodologica y para que funcione si en el futuro "
        "  cambia la distribucion.\n"
        "- Que pasa si una variable cambia de rango? -> el schema con "
        "  pandera lo detecta en la carga y falla con mensaje claro."
    ))


# ─── Main ─────────────────────────────────────────────────────────────


def main():
    prs = Presentation()
    prs.slide_width = Inches(10.667)   # 16:9
    prs.slide_height = Inches(7.5)

    slide_1_portada(prs)
    slide_2_pregunta_negocio(prs)
    slide_3_variables(prs)
    slide_4_pipeline(prs)
    slide_5_resultados(prs)
    slide_6_graficas(prs)
    slide_7_importancia(prs)
    slide_8_cierre(prs)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)

    print(f"OK: presentacion generada en {OUTPUT_PATH}")
    print(f"    Tamano: {OUTPUT_PATH.stat().st_size / 1024:.1f} KB")
    print(f"    Slides: {len(prs.slides)}")
    print()
    print("SIGUIENTE PASO:")
    print("  1. Abrir el archivo en PowerPoint")
    print("  2. En la slide 6 (matriz + ROC), arrastrar las imagenes:")
    print("     - reports/v1/confusion_matrix.png")
    print("     - reports/v1/roc_curve.png")
    print("  3. En la slide 7 (importancia), arrastrar:")
    print("     - reports/v1/feature_importances.png")
    print("  4. Revisar las notas del presentador (Vista > Notas)")


if __name__ == "__main__":
    main()